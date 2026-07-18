#!/usr/bin/env python3
"""Export an HTML Report Stable Base file according to its declared mode."""

from __future__ import annotations

import argparse
import base64
import json
import mimetypes
import re
import tempfile
from html import unescape
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
from playwright.sync_api import sync_playwright


SLIDE_W = 1440
SLIDE_H = 810
CHROME_CANDIDATES = [
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/Applications/Chromium.app/Contents/MacOS/Chromium",
]

CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}

SHAPE_TYPES = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "roundedRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "line": MSO_SHAPE.RECTANGLE,
}

DEFAULT_FONT_FAMILY = "Microsoft YaHei"
PPT_LEVEL_FONT_SIZES = {
    "h1": 60,
    "h2": 28,
    "h3": 20,
    "metric": 30,
    "body": 14,
}
PPT_ROLE_FONT_SIZES = {
    "header": 12,
    "footer": 12,
    "pageNumber": 12,
}


def data_uri(path: Path) -> str:
    mime = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    if path.suffix.lower() == ".svg":
        mime = "image/svg+xml"
    if path.suffix.lower() == ".webp":
        mime = "image/webp"
    return f"data:{mime};base64,{base64.b64encode(path.read_bytes()).decode('ascii')}"


def inline_local_refs(html_path: Path, out_path: Path) -> None:
    root = html_path.parent
    html = html_path.read_text(encoding="utf-8")

    def replace_script(match: re.Match[str]) -> str:
        src = match.group("src")
        if src.startswith(("http://", "https://", "data:")):
            return match.group(0)
        script_path = (root / src).resolve()
        if not script_path.exists():
            return match.group(0)
        return f'<script data-inline-asset="{src}">{script_path.read_text(encoding="utf-8")}</script>'

    html = re.sub(
        r'<script\b(?P<attrs>[^>]*)\bsrc=["\'](?P<src>[^"\']+)["\'][^>]*>\s*</script>',
        replace_script,
        html,
        flags=re.I,
    )

    def inline_css_urls(css: str, base_dir: Path) -> str:
        def replace_url(match: re.Match[str]) -> str:
            ref = match.group("ref").strip().strip("\"'")
            if ref.startswith(("http://", "https://", "data:", "#")):
                return match.group(0)
            path = (base_dir / ref).resolve()
            return f"url('{data_uri(path)}')" if path.exists() and path.is_file() else match.group(0)

        return re.sub(r"url\((?P<ref>[^)]+)\)", replace_url, css, flags=re.I)

    def replace_stylesheet(match: re.Match[str]) -> str:
        href = match.group("href")
        if href.startswith(("http://", "https://", "data:")):
            return match.group(0)
        css_path = (root / href).resolve()
        if not css_path.exists():
            return match.group(0)
        css = inline_css_urls(css_path.read_text(encoding="utf-8"), css_path.parent)
        return f'<style data-inline-asset="{href}">{css}</style>'

    html = re.sub(
        r'<link\b(?=[^>]*\brel=["\']stylesheet["\'])[^>]*\bhref=["\'](?P<href>[^"\']+)["\'][^>]*>',
        replace_stylesheet,
        html,
        flags=re.I,
    )

    def replace_image(match: re.Match[str]) -> str:
        src = match.group("src")
        if src.startswith(("http://", "https://", "data:")):
            return match.group(0)
        image_path = (root / src).resolve()
        if not image_path.exists() or not image_path.is_file():
            return match.group(0)
        return match.group(0).replace(src, data_uri(image_path), 1)

    html = re.sub(
        r'<img\b[^>]*\bsrc=["\'](?P<src>[^"\']+)["\'][^>]*>',
        replace_image,
        html,
        flags=re.I,
    )
    html = sanitize_standalone_html(html)

    asset_refs = sorted(
        set(re.findall(r'(?<!data:)(?:assets|data|images|media)/[A-Za-z0-9_./\-\u4e00-\u9fff]+(?:\.(?:svg|png|jpg|jpeg|webp|ico|json|js))', html)),
        key=len,
        reverse=True,
    )
    for ref in asset_refs:
        path = (root / ref).resolve()
        if path.exists() and path.is_file():
            if path.suffix.lower() in {".png", ".jpg", ".jpeg", ".svg", ".webp", ".ico"}:
                html = html.replace(ref, data_uri(path))
            elif path.suffix.lower() in {".json"}:
                html = html.replace(ref, data_uri(path))

    html = html.replace(
        "<head>",
        "<head>\n  <!-- Standalone single-file export: local assets and scripts are inlined. -->",
        1,
    )
    out_path.write_text(html, encoding="utf-8")


def remove_class_tokens(match: re.Match[str], tokens: set[str]) -> str:
    tag = match.group(0)
    class_match = re.search(r'class=["\']([^"\']*)["\']', tag)
    if not class_match:
        return tag
    classes = [part for part in class_match.group(1).split() if part not in tokens]
    replacement = f'class="{" ".join(classes)}"'
    return tag[: class_match.start()] + replacement + tag[class_match.end() :]


def add_class_token_to_tag(tag_or_match: str | re.Match[str], token: str) -> str:
    tag = tag_or_match if isinstance(tag_or_match, str) else tag_or_match.group(0)
    class_match = re.search(r'class=["\']([^"\']*)["\']', tag)
    if not class_match:
        return tag[:-1] + f' class="{token}">'
    classes = class_match.group(1).split()
    if token not in classes:
        classes.append(token)
    replacement = f'class="{" ".join(classes)}"'
    return tag[: class_match.start()] + replacement + tag[class_match.end() :]


def add_boolean_attr_to_tag(match: re.Match[str], attr: str) -> str:
    tag = match.group(0)
    if re.search(rf"\b{re.escape(attr)}(?:\b|=)", tag, flags=re.I):
        return tag
    return tag[:-1] + f" {attr}>"


def hide_standalone_drawer_controls(html: str) -> str:
    style = (
        "<style data-standalone-export-ui>"
        "body.standalone-export .drawer-toggle,"
        "body.standalone-export .style-drawer{display:none!important;}"
        "</style>"
    )
    if "data-standalone-export-ui" not in html:
        if "</head>" in html:
            html = html.replace("</head>", f"{style}</head>", 1)
        else:
            html = style + html
    html = re.sub(
        r'<button\b[^>]*\bid=["\']drawerToggle["\'][^>]*>',
        lambda match: add_boolean_attr_to_tag(match, "hidden"),
        html,
        count=1,
        flags=re.I,
    )
    html = re.sub(
        r'<aside\b[^>]*\bid=["\']styleDrawer["\'][^>]*>',
        lambda match: add_boolean_attr_to_tag(match, "hidden"),
        html,
        count=1,
        flags=re.I,
    )
    return html


def standalone_nav_style(html: str) -> str:
    body_match = re.search(r"<body\b([^>]*)>", html, flags=re.I)
    classes = body_match.group(1) if body_match else ""
    if re.search(r'class=["\'][^"\']*\bnav-hidden\b', classes):
        return "hidden"
    if re.search(r'class=["\'][^"\']*\bnav-dots\b', classes):
        return "dots"
    return "numbered"


def sanitize_standalone_html(html: str) -> str:
    html_close = html.lower().find("</html>")
    if html_close >= 0:
        html = html[: html_close + len("</html>")]
    nav_style = standalone_nav_style(html)
    html = re.sub(
        r'(<nav\b[^>]*\bid=["\']nav["\'][^>]*>)([\s\S]*?)(</nav>)',
        r"\1\3",
        html,
        flags=re.I,
    )
    html = re.sub(
        r'(<div\b[^>]*\bid=["\']typographyLabelLayer["\'][^>]*>)([\s\S]*?)(</div>)',
        r"\1\3",
        html,
        flags=re.I,
    )
    html = re.sub(
        r'<body\b[^>]*>',
        lambda match: add_class_token_to_tag(remove_class_tokens(match, {"drawer-open", "edit-mode"}), "standalone-export"),
        html,
        count=1,
        flags=re.I,
    )
    html = re.sub(
        r'<aside\b[^>]*\bid=["\']styleDrawer["\'][^>]*>',
        lambda match: remove_class_tokens(match, {"open"}),
        html,
        count=1,
        flags=re.I,
    )
    html = hide_standalone_drawer_controls(html)
    html = re.sub(r'\scontenteditable(?:=["\'][^"\']*["\'])?', "", html, flags=re.I)
    html = re.sub(r'\sdata-typography-label=["\'][^"\']*["\']', "", html, flags=re.I)
    for token in ("selected-box", "selected-col"):
        html = re.sub(
            r'class=["\']([^"\']*)["\']',
            lambda match, token=token: f'class="{" ".join(part for part in match.group(1).split() if part != token)}"',
            html,
        )
    html = re.sub(
        r"^(\s*)setNavStyle\('numbered'\);$",
        rf"\1setNavStyle('{nav_style}');",
        html,
        count=1,
        flags=re.M,
    )
    html = html.replace(
        "<script>setNavStyle('numbered');</script>",
        f"<script>setNavStyle('{nav_style}');</script>",
        1,
    )
    return html


def parse_color(value: str | None, fallback: str = "#000000") -> RGBColor:
    color = (value or fallback).strip()
    if color == "none":
        color = fallback
    rgb_match = re.match(r"rgba?\((\d+),\s*(\d+),\s*(\d+)", color)
    if rgb_match:
        return RGBColor(*(int(part) for part in rgb_match.groups()))
    if color.startswith("#"):
        color = color[1:]
    if len(color) == 3:
        color = "".join(ch * 2 for ch in color)
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def set_fill(shape, fill: str | None) -> None:
    if not fill or fill == "none":
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = parse_color(fill, "#ffffff")


def set_line(shape, line: str | dict | None) -> None:
    if line == "none":
        shape.line.fill.background()
        return
    if isinstance(line, dict):
        fill = line.get("fill") or line.get("color") or "#dbe4f0"
        width = line.get("width")
    else:
        fill = line or "#dbe4f0"
        width = None
    shape.line.color.rgb = parse_color(fill, "#dbe4f0")
    if width:
        shape.line.width = Inches(float(width) / 96)


def concrete_font_family(font_family: str | None, fallback: str = DEFAULT_FONT_FAMILY) -> str:
    if not font_family:
        return fallback
    for part in font_family.split(","):
        name = part.strip().strip("\"'")
        if name and name.lower() not in {"serif", "sans-serif", "monospace", "cursive", "fantasy", "system-ui"}:
            return name
    return fallback


def set_run_font(run, font_family: str | None, font_size_px: float, bold: bool, color: str | None) -> None:
    family = concrete_font_family(font_family)
    run.font.name = family
    run.font.size = Pt(float(font_size_px))
    run.font.bold = bold
    run.font.color.rgb = parse_color(color, "#1f2937")
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        element = rpr.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            rpr.append(element)
        element.set("typeface", family)


def ppt_font_size(obj: dict, fallback: float = 14) -> float:
    role = obj.get("pptRole")
    if role in PPT_ROLE_FONT_SIZES:
        return PPT_ROLE_FONT_SIZES[role]
    level = obj.get("pptLevel")
    if level == "metric" and obj.get("fontSize") is not None:
        return float(obj["fontSize"])
    if level in PPT_LEVEL_FONT_SIZES:
        return PPT_LEVEL_FONT_SIZES[level]
    return float(obj.get("fontSize", fallback))


def extract_deck_json(html_path: Path, explicit_json: str | None = None, executable: str | None = None) -> dict:
    if explicit_json:
        return json.loads(Path(explicit_json).read_text(encoding="utf-8"))
    html = html_path.read_text(encoding="utf-8")
    match = re.search(
        r'<script[^>]+id=["\']html-pptx-data["\'][^>]*>(.*?)</script>',
        html,
        flags=re.S | re.I,
    )
    if not match:
        raise RuntimeError(
            "Missing html-pptx-data JSON contract. This PPTX exporter only handles "
            "PPT-mode HTML generated by html-report-stable-base with structured JSON."
        )
    if "syncDeckJsonFromDom" in html:
        with sync_playwright() as p:
            kwargs = {"headless": True}
            if executable:
                kwargs["executable_path"] = executable
            browser = p.chromium.launch(**kwargs)
            page = browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H}, device_scale_factor=1)
            page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
            page.wait_for_timeout(600)
            deck = page.evaluate("window.syncDeckJsonFromDom ? window.syncDeckJsonFromDom() : JSON.parse(document.getElementById('html-pptx-data').textContent)")
            browser.close()
            return deck
    return json.loads(unescape(match.group(1)).strip())


def px_to_emu(value: float, total_px: float, total_emu: int) -> int:
    return int(float(value) / float(total_px) * total_emu)


def object_position(obj: dict, meta: dict, prs: Presentation) -> tuple[int, int, int, int]:
    width = float(meta.get("width", SLIDE_W))
    height = float(meta.get("height", SLIDE_H))
    return (
        px_to_emu(obj.get("x", 0), width, prs.slide_width),
        px_to_emu(obj.get("y", 0), height, prs.slide_height),
        px_to_emu(obj.get("w", obj.get("width", 0)), width, prs.slide_width),
        px_to_emu(obj.get("h", obj.get("height", 0)), height, prs.slide_height),
    )


def add_text(slide, obj: dict, meta: dict, prs: Presentation) -> None:
    left, top, width, height = object_position(obj, meta, prs)
    box = slide.shapes.add_textbox(left, top, width, height)
    box.name = obj.get("name", "text")
    tf = box.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.text = str(obj.get("text", ""))
    align = obj.get("align")
    if align == "center":
        paragraph.alignment = PP_ALIGN.CENTER
    elif align == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    elif align == "left":
        paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    set_run_font(
        run,
        obj.get("fontFamily") or meta.get("fontFamily") or DEFAULT_FONT_FAMILY,
        ppt_font_size(obj, 14),
        bool(obj.get("bold", False)),
        obj.get("color"),
    )


def add_shape(slide, obj: dict, meta: dict, prs: Presentation) -> None:
    left, top, width, height = object_position(obj, meta, prs)
    shape_type = SHAPE_TYPES.get(obj.get("shape", "rect"), MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.name = obj.get("name", "shape")
    set_fill(shape, obj.get("fill", "#ffffff"))
    line = obj.get("line", "#dbe4f0")
    set_line(shape, line if line == "none" else {"color": line, "width": obj.get("lineWidth")})
    if obj.get("radius") is not None and obj.get("shape") in {"roundRect", "roundedRect"}:
        try:
            radius = float(obj.get("radius", 8))
            min_side = max(1.0, min(float(obj.get("w", 1)), float(obj.get("h", 1))))
            shape.adjustments[0] = max(0.0, min(0.18, radius / min_side))
        except Exception:
            pass


def add_table(slide, obj: dict, meta: dict, prs: Presentation) -> None:
    columns = obj.get("columns") or []
    rows = obj.get("rows") or []
    if not columns:
        raise RuntimeError(f"Table {obj.get('name', '')} has no columns")
    left, top, width, height = object_position(obj, meta, prs)
    table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width, height)
    table_shape.name = obj.get("name", "table")
    table = table_shape.table
    column_widths = obj.get("columnWidths") or []
    if column_widths:
        total = sum(float(x) for x in column_widths) or 1
        for i, value in enumerate(column_widths[: len(columns)]):
            table.columns[i].width = int(width * float(value) / total)
    for c, heading in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = str(heading)
        cell.fill.solid()
        cell.fill.fore_color.rgb = parse_color(obj.get("headerFill"), "#172033")
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                set_run_font(run, meta.get("fontFamily") or DEFAULT_FONT_FAMILY, 14, True, "#ffffff")
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row[: len(columns)]):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = parse_color("#f8fafc" if r % 2 == 0 else "#ffffff")
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_run_font(run, meta.get("fontFamily") or DEFAULT_FONT_FAMILY, 14, c == 0, "#334155")


def apply_chart_series_colors(chart, series_configs: list[dict]) -> None:
    for index, series_config in enumerate(series_configs):
        if index >= len(chart.series):
            continue
        color = series_config.get("color")
        if not color:
            continue
        series = chart.series[index]
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = parse_color(color)
        except Exception:
            pass
        try:
            series.format.line.color.rgb = parse_color(color)
            series.format.line.width = Pt(2.25)
        except Exception:
            pass


def add_single_chart(slide, obj: dict, meta: dict, prs: Presentation, chart_type_name: str, series_configs: list[dict], name_suffix: str = "", overlay: bool = False):
    chart_type = CHART_TYPES.get(chart_type_name)
    if not chart_type:
        raise RuntimeError(f"Unsupported chart type: {chart_type_name}")
    left, top, width, height = object_position(obj, meta, prs)
    data = CategoryChartData()
    data.categories = obj.get("categories") or []
    if chart_type_name == "pie":
        values = obj.get("values")
        if values is None and series_configs:
            values = series_configs[0].get("values", [])
        data.add_series(obj.get("seriesName", series_configs[0].get("name", "Series") if series_configs else "Series"), values or [])
    else:
        for series in series_configs:
            data.add_series(series.get("name", "Series"), series.get("values", []))
    frame = slide.shapes.add_chart(chart_type, left, top, width, height, data)
    frame.name = obj.get("name", "chart") + name_suffix
    chart = frame.chart
    chart.has_legend = bool(obj.get("legend", True))
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM if chart_type_name != "pie" else XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    apply_chart_series_colors(chart, series_configs)
    if overlay:
        chart.has_legend = False
        try:
            chart.category_axis.visible = False
            chart.value_axis.visible = False
        except Exception:
            pass
        try:
            chart.chart_area.format.fill.background()
            chart.plot_area.format.fill.background()
        except Exception:
            pass
    return frame


def add_chart(slide, obj: dict, meta: dict, prs: Presentation) -> None:
    chart_type_name = obj.get("chartType") or obj.get("type")
    series_configs = obj.get("series") or []
    if chart_type_name not in CHART_TYPES:
        left, top, width, height = object_position(obj, meta, prs)
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        placeholder.name = obj.get("name", "chart")
        set_fill(placeholder, "#f8fafc")
        set_line(placeholder, {"color": "#cbd5e1", "width": 1})
        placeholder.text_frame.text = "图表未嵌入"
        for paragraph in placeholder.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                set_run_font(run, meta.get("fontFamily") or DEFAULT_FONT_FAMILY, 12, False, "#64748b")
        return
    series_types = {series.get("type", chart_type_name) for series in series_configs}
    if "bar" in series_types and "line" in series_types:
        bar_series = [series for series in series_configs if series.get("type", chart_type_name) == "bar"]
        line_series = [series for series in series_configs if series.get("type", chart_type_name) == "line"]
        add_single_chart(slide, {**obj, "legend": False}, meta, prs, "bar", bar_series, "-bar")
        add_single_chart(slide, {**obj, "legend": False}, meta, prs, "line", line_series, "-line", overlay=True)
        return
    add_single_chart(slide, obj, meta, prs, chart_type_name, series_configs)


def add_image(slide, obj: dict, meta: dict, prs: Presentation, root: Path) -> None:
    src = obj.get("src")
    if not src:
        return
    left, top, width, height = object_position(obj, meta, prs)
    if str(src).startswith("data:image/"):
        header, encoded = str(src).split(",", 1)
        suffix = ".png" if "png" in header else ".jpg"
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as handle:
            handle.write(base64.b64decode(encoded))
            image_path = Path(handle.name)
        try:
            slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        finally:
            image_path.unlink(missing_ok=True)
        return
    image_path = (root / src).resolve()
    if not image_path.exists():
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        placeholder.name = obj.get("name", "image")
        set_fill(placeholder, "#edf2f7")
        set_line(placeholder, {"color": "#cbd5e1", "width": 1})
        placeholder.text_frame.text = "图片未嵌入"
        for paragraph in placeholder.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                set_run_font(run, meta.get("fontFamily") or DEFAULT_FONT_FAMILY, 12, False, "#64748b")
        return
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def export_editable_pptx(deck: dict, out_path: Path, root: Path) -> None:
    meta = deck.get("meta") or {}
    if not deck.get("slides"):
        raise RuntimeError("html-pptx-data has no slides")
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for slide_data in deck["slides"]:
        slide = prs.slides.add_slide(blank)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = parse_color(slide_data.get("background"), "#ffffff")
        for obj in slide_data.get("objects") or []:
            kind = obj.get("type")
            if kind == "text":
                add_text(slide, obj, meta, prs)
            elif kind == "shape":
                add_shape(slide, obj, meta, prs)
            elif kind == "table":
                add_table(slide, obj, meta, prs)
            elif kind == "chart":
                add_chart(slide, obj, meta, prs)
            elif kind == "image":
                add_image(slide, obj, meta, prs, root)
            else:
                raise RuntimeError(f"Unsupported PPTX object type: {kind}")
    prs.save(out_path)


def chrome_path(explicit: str | None) -> str | None:
    if explicit:
        return explicit
    for candidate in CHROME_CANDIDATES:
        if Path(candidate).exists():
            return candidate
    return None


def capture_slides(html_path: Path, shot_dir: Path, executable: str | None) -> list[Path]:
    with sync_playwright() as p:
        kwargs = {"headless": True}
        if executable:
            kwargs["executable_path"] = executable
        browser = p.chromium.launch(**kwargs)
        page = browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H}, device_scale_factor=2)
        page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
        page.wait_for_timeout(600)
        slide_count = page.evaluate("document.querySelectorAll('.slide').length")
        shots: list[Path] = []
        for i in range(slide_count):
            page.evaluate("(i) => window.goTo ? window.goTo(i) : document.querySelectorAll('.slide')[i].scrollIntoView()", i)
            page.wait_for_timeout(350)
            path = shot_dir / f"slide_{i + 1:02d}.png"
            page.locator(".slide.active-slide").screenshot(path=str(path))
            shots.append(path)
        browser.close()
    return shots


def export_pdf(shots: list[Path], out_path: Path) -> None:
    images = [Image.open(p).convert("RGB") for p in shots]
    if not images:
        raise RuntimeError("No slide screenshots captured")
    images[0].save(out_path, save_all=True, append_images=images[1:])


def export_screenshot_pptx(shots: list[Path], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for shot in shots:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(shot), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out_path)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("html", help="Source index.html")
    parser.add_argument("--out-dir", default=None, help="Export directory")
    parser.add_argument("--formats", default=None, help="Comma list: html,pdf,pptx")
    parser.add_argument("--chrome-path", default=None, help="Optional Chrome/Chromium executable")
    parser.add_argument("--deck-json", default=None, help="Optional structured JSON file for editable PPTX")
    args = parser.parse_args()

    html_path = Path(args.html).resolve()
    out_dir = Path(args.out_dir).resolve() if args.out_dir else html_path.parent / "exports"
    out_dir.mkdir(parents=True, exist_ok=True)
    source_html = html_path.read_text(encoding="utf-8")
    mode_match = re.search(r'<body[^>]+data-report-mode=["\']([^"\']+)', source_html, flags=re.I)
    if not mode_match or mode_match.group(1) not in {"single-page", "ppt"}:
        raise SystemExit("HTML must declare data-report-mode as single-page or ppt")
    mode = mode_match.group(1)
    requested = args.formats or ("html" if mode == "single-page" else "html,pdf,pptx")
    formats = {x.strip().lower() for x in requested.split(",") if x.strip()}
    unsupported = formats - {"html", "pdf", "pptx"}
    if unsupported:
        raise SystemExit(f"Unsupported formats: {', '.join(sorted(unsupported))}")
    if mode == "single-page" and formats - {"html"}:
        raise SystemExit("single-page mode only supports standalone HTML export")

    outputs: list[Path] = []
    if "html" in formats:
        single = out_dir / f"{html_path.stem}_single.html"
        inline_local_refs(html_path, single)
        outputs.append(single)

    if "pdf" in formats:
        with tempfile.TemporaryDirectory() as tmp:
            shot_dir = Path(tmp)
            shots = capture_slides(html_path, shot_dir, chrome_path(args.chrome_path))
            pdf = out_dir / f"{html_path.stem}.pdf"
            export_pdf(shots, pdf)
            outputs.append(pdf)

    if "pptx" in formats:
        pptx = out_dir / f"{html_path.stem}.pptx"
        deck = extract_deck_json(html_path, args.deck_json, chrome_path(args.chrome_path))
        export_editable_pptx(deck, pptx, html_path.parent)
        outputs.append(pptx)

    for output in outputs:
        print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
