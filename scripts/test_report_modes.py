#!/usr/bin/env python3
"""Regression checks for HTML Report modes and editor contract."""

from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
BUILDER = ROOT / "scripts" / "build_html_report_from_model.py"
CHECKER = ROOT / "scripts" / "check_html_report.py"
EXPORTER = ROOT / "scripts" / "export_html_report.py"
SINGLE_TEMPLATE = ROOT / "assets" / "template" / "html-report-single-base.html"
PPT_TEMPLATE = ROOT / "assets" / "template" / "html-report-ppt-base.html"
EDITOR_JS = ROOT / "assets" / "template" / "shared" / "report-editor.js"
EDITOR_CSS = ROOT / "assets" / "template" / "shared" / "report-editor.css"


def run_builder(model: dict, template: Path, output: Path) -> subprocess.CompletedProcess[str]:
    model_path = output.with_suffix(".json")
    model_path.write_text(json.dumps(model, ensure_ascii=False), encoding="utf-8")
    return subprocess.run(
        [sys.executable, str(BUILDER), str(model_path), "--template", str(template), "--out", str(output)],
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )


def assert_true(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def test_mode_is_required() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        result = run_builder(
            {"title": "Missing mode", "slides": [{"title": "One", "html": "<p>One</p>", "objects": []}]},
            PPT_TEMPLATE,
            Path(tmp) / "index.html",
        )
        assert_true(result.returncode != 0, "builder must reject a model without mode")
        message = result.stdout + result.stderr
        assert_true("can't open file" not in message, "builder entrypoint must exist before mode validation can be tested")
        assert_true("model.mode" in message.lower(), "missing-mode error must mention model.mode confirmation")


def test_single_page_contract() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        output = Path(tmp) / "index.html"
        result = run_builder(
            {
                "mode": "single-page",
                "title": "Long report",
                "sections": [
                    {"id": "intro", "title": "超过六个字的章节标题", "html": '<h1 data-editable-element="text">Intro</h1>'},
                    {"id": "data", "title": "数据", "html": '<p data-editable-element="text">Body</p>'},
                ],
            },
            SINGLE_TEMPLATE,
            output,
        )
        assert_true(result.returncode == 0, result.stderr or result.stdout)
        html = output.read_text(encoding="utf-8")
        assert_true('data-report-mode="single-page"' in html, "single-page output must declare its mode")
        assert_true("html-pptx-data" not in html, "single-page output must not contain PPTX JSON")
        assert_true('id="exportPpt"' not in html and 'id="printPdf"' not in html, "single-page output must not expose PPT/PDF controls")
        assert_true('id="nav"' not in html, "single-page output must not contain PPT navigation")
        assert_true('id="sectionNav"' in html, "single-page output must expose section navigation")
        assert_true('id="sectionNavStyleControl"' in html, "single-page drawer must expose section navigation controls")
        for style in ("hidden", "left", "right"):
            assert_true(f'data-section-nav-style="{style}"' in html, f"single-page navigation must support {style}")
        assert_true('class="single-page-mode section-nav-right"' in html, "single-page navigation must default to right")
        assert_true('id="downloadHtml"' in html, "single-page output must keep standalone HTML export")
        assert_true('<script src="shared/report-editor.js"' not in html, "generated output must inline the shared editor")
        checked = subprocess.run([sys.executable, str(CHECKER), str(output)], text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
        assert_true(checked.returncode == 0, checked.stderr or checked.stdout)
        rejected = subprocess.run([sys.executable, str(EXPORTER), str(output), "--formats", "pptx"], text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
        assert_true(rejected.returncode != 0 and "only supports" in (rejected.stdout + rejected.stderr), "single-page mode must reject PPTX export")


def test_ppt_contract_remains_available() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        output = Path(tmp) / "index.html"
        result = run_builder(
            {
                "mode": "ppt",
                "title": "Deck",
                "slides": [{"title": "One", "html": '<h1 data-pptx-name="title">One</h1>', "objects": []}],
            },
            PPT_TEMPLATE,
            output,
        )
        assert_true(result.returncode == 0, result.stderr or result.stdout)
        html = output.read_text(encoding="utf-8")
        assert_true('data-report-mode="ppt"' in html, "PPT output must declare its mode")
        assert_true("html-pptx-data" in html, "PPT output must keep PPTX JSON")
        assert_true('id="nav"' in html and 'id="exportPpt"' in html, "PPT output must keep navigation and PPT export")
        checked = subprocess.run([sys.executable, str(CHECKER), str(output)], text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
        assert_true(checked.returncode == 0, checked.stderr or checked.stdout)


def test_ppt_export_syncs_missing_named_dom_objects() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        output = tmp_path / "index.html"
        result = run_builder(
            {
                "mode": "ppt",
                "title": "Deck",
                "slides": [
                    {
                        "title": "One",
                        "html": """
<h2 class="editable editable-box" data-editable-element="text" data-pptx-name="missing-title">Missing title should export</h2>
<div class="panel editable-box" data-editable-element="shape" data-pptx-name="missing-card">
  <h3 class="editable" data-editable-element="text" data-pptx-name="missing-heading">Missing heading should export</h3>
  <p class="editable" data-editable-element="text" data-pptx-name="missing-body">Missing body should export</p>
</div>
""",
                        "objects": [],
                    }
                ],
            },
            PPT_TEMPLATE,
            output,
        )
        assert_true(result.returncode == 0, result.stderr or result.stdout)
        exported = subprocess.run(
            [
                sys.executable,
                str(EXPORTER),
                str(output),
                "--out-dir",
                str(tmp_path / "exports"),
                "--formats",
                "pptx",
            ],
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=False,
        )
        assert_true(exported.returncode == 0, exported.stderr or exported.stdout)
        pptx_path = tmp_path / "exports" / "index.pptx"
        prs = Presentation(pptx_path)
        exported_text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        assert_true("Missing title should export" in exported_text, "PPTX export must include missing named title text")
        assert_true("Missing heading should export" in exported_text, "PPTX export must include nested missing heading text")
        assert_true("Missing body should export" in exported_text, "PPTX export must include nested missing body text")


def test_ppt_export_syncs_unnamed_editable_dom_objects() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        output = tmp_path / "index.html"
        result = run_builder(
            {
                "mode": "ppt",
                "title": "Deck",
                "slides": [
                    {
                        "title": "One",
                        "html": """
<h2 class="editable editable-box" data-editable-element="text">Unnamed title should export</h2>
<div class="panel editable-box" data-editable-element="shape">
  <h3 class="editable" data-editable-element="text">Unnamed heading should export</h3>
  <ul>
    <li class="editable" data-editable-element="text">Unnamed bullet should export</li>
  </ul>
</div>
""",
                        "objects": [],
                    }
                ],
            },
            PPT_TEMPLATE,
            output,
        )
        assert_true(result.returncode == 0, result.stderr or result.stdout)
        exported = subprocess.run(
            [
                sys.executable,
                str(EXPORTER),
                str(output),
                "--out-dir",
                str(tmp_path / "exports"),
                "--formats",
                "pptx",
            ],
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=False,
        )
        assert_true(exported.returncode == 0, exported.stderr or exported.stdout)
        pptx_path = tmp_path / "exports" / "index.pptx"
        prs = Presentation(pptx_path)
        exported_text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        assert_true("Unnamed title should export" in exported_text, "PPTX export must include unnamed title text")
        assert_true("Unnamed heading should export" in exported_text, "PPTX export must include unnamed heading text")
        assert_true("Unnamed bullet should export" in exported_text, "PPTX export must include unnamed bullet text")


def test_shared_editor_contract() -> None:
    source = EDITOR_JS.read_text(encoding="utf-8")
    for element_type in ("text", "shape", "chart", "image"):
        assert_true(element_type in source, f"shared editor must support {element_type}")
    for control in (
        "elementWidth",
        "elementHeight",
        "elementFontSize",
        "elementColor",
        "elementFontWeight",
        "elementTextAlign",
        "elementLineHeight",
        "elementBackground",
        "elementBorderColor",
        "elementBorderWidth",
        "elementBorderRadius",
    ):
        assert_true(control in source, f"shared editor must bind {control}")
    assert_true("selectedElement" in source, "editor must maintain exactly one current element selection")
    assert_true("getInstanceByDom" in source and ".resize()" in source, "chart size changes must resize ECharts")
    assert_true("currentStageScale" in source, "PPT element dimensions must convert through the current stage scale")
    for contract in ("setSectionNavStyle", "buildSectionNav", "shortSectionTitle", "typographyLevel", "refreshTypographyLabels"):
        assert_true(contract in source, f"shared editor must expose {contract}")
    assert_true("Array.from(String(title" in source and ".slice(0, 6)" in source, "single-page section nav labels must truncate by Unicode characters")
    for style in ("hidden", "left", "right"):
        assert_true(f"'{style}'" in source, f"single-page nav logic must support {style}")


def test_shared_accordion_contract() -> None:
    source = EDITOR_JS.read_text(encoding="utf-8")
    styles = EDITOR_CSS.read_text(encoding="utf-8")
    for contract in (
        "initAccordionGroups",
        "control-group-toggle",
        "control-group-content",
        "aria-expanded",
    ):
        assert_true(contract in source, f"shared editor must expose accordion contract: {contract}")
    assert_true("setAttribute('aria-expanded', 'false')" in source, "accordion groups must initialize collapsed")
    assert_true("content.hidden = true" in source, "accordion content must be hidden by default")
    assert_true("document.querySelectorAll('.drawer-body > .control-group').forEach" in source, "accordion opening must query every current group")
    assert_true("setAccordionGroupOpen(otherGroup, false)" in source, "accordion opening must close other groups")
    assert_true("globalControlInitialized" in source, "global control listeners must be idempotent across editor initialization")
    assert_true("elementControlInitialized" in source and "documentClickBound" in source, "element and document listeners must be idempotent across editor initialization")
    assert_true("group.hidden" in source, "accordion must preserve element-control-group hidden state")
    assert_true(".control-group[hidden]" in styles, "hidden element control groups must remain hidden")
    assert_true(".control-group-content[hidden]" in styles, "collapsed accordion content must remain hidden")
    for template in (SINGLE_TEMPLATE, PPT_TEMPLATE):
        html = template.read_text(encoding="utf-8")
        assert_true('href="shared/report-editor.css"' in html, f"{template.name} must use shared editor CSS")
        assert_true('src="shared/report-editor.js"' in html, f"{template.name} must use shared editor JS")
        assert_true('class="drawer-body"' in html and 'class="control-group"' in html, f"{template.name} must expose drawer groups")


def test_ppt_drawer_accordion_groups_do_not_stretch() -> None:
    ppt = PPT_TEMPLATE.read_text(encoding="utf-8")
    assert_true("align-content: start" in ppt, "PPT drawer body must not stretch collapsed accordion groups")
    assert_true("grid-auto-rows: max-content" in ppt, "PPT drawer rows must size to content")
    assert_true("align-self: start" in ppt, "PPT drawer control groups must stay compact")


def test_metric_typography_is_a_seventh_global_level() -> None:
    single = SINGLE_TEMPLATE.read_text(encoding="utf-8")
    ppt = PPT_TEMPLATE.read_text(encoding="utf-8")
    for name, html in (("single-page", single), ("ppt", ppt)):
        assert_true("--metric-size:" in html, f"{name} template must define the metric typography level")
        assert_true('data-global-var="--metric-size"' in html, f"{name} drawer must expose metric size")
        assert_true("指标字号" in html, f"{name} drawer must name the metric size control")
        assert_true(".metric .value" in html and "font-size: var(--metric-size)" in html, f"{name} metric values must use metric size")
        assert_true(".metric .label" in html and "font-size: var(--body-size)" in html, f"{name} metric labels must remain body text")
    for level in ("'L1'", "'L2'", "'L3'", "'L4'", "'指标'", "'正文'", "'备注'"):
        assert_true(level in EDITOR_JS.read_text(encoding="utf-8"), f"single-page typography labels must include {level}")
    assert_true(".typography-label" in single, "single-page template must include visible typography labels")
    assert_true("element.matches('[data-ppt-level=\"metric\"], .metric .value, .metric strong')" in ppt, "PPT typography classification must only recognize metric values")
    assert_true("return '指标'" in ppt, "PPT edit labels must display the metric typography level")
    assert_true("obj.pptLevel = 'metric'" in ppt, "PPT DOM sync must export metric-card text as metric typography")


def main() -> int:
    tests = [
        test_mode_is_required,
        test_single_page_contract,
        test_ppt_contract_remains_available,
        test_ppt_export_syncs_missing_named_dom_objects,
        test_ppt_export_syncs_unnamed_editable_dom_objects,
        test_shared_editor_contract,
        test_shared_accordion_contract,
        test_ppt_drawer_accordion_groups_do_not_stretch,
        test_metric_typography_is_a_seventh_global_level,
    ]
    for test in tests:
        test()
        print(f"[OK] {test.__name__}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
